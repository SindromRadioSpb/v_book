"""PPTX extractor."""
import logging
from pathlib import Path

logger = logging.getLogger(__name__)


def extract_text(file_path: Path) -> str:
    """
    Extract text from PPTX file.

    Extracts all text from slides, including:
    - Slide titles
    - Body text
    - Text in shapes
    - Table cells

    Note: Does not extract speaker notes or comments.
    Returns empty string if python-pptx is not available.
    """
    logger.info(f"Extracting text from PPTX: {file_path}")

    try:
        from pptx import Presentation
    except ImportError:
        logger.error("python-pptx not installed. Install with: pip install python-pptx")
        raise ImportError("python-pptx library is required for PPTX extraction")

    try:
        prs = Presentation(file_path)
        text_runs = []

        # Iterate through all slides
        for slide_idx, slide in enumerate(prs.slides, start=1):
            # Extract text from all shapes in the slide
            for shape in slide.shapes:
                # Check if shape has text frame
                if hasattr(shape, "text") and shape.text:
                    text = shape.text.strip()
                    if text:
                        text_runs.append(text)

                # Extract text from tables
                if hasattr(shape, "table"):
                    table = shape.table
                    for row in table.rows:
                        for cell in row.cells:
                            text = cell.text.strip()
                            if text:
                                text_runs.append(text)

        result = "\n".join(text_runs)
        logger.info(f"Extracted {len(result)} characters from PPTX ({len(prs.slides)} slides)")
        return result

    except Exception as e:
        logger.exception(f"Failed to extract text from PPTX: {file_path}")
        raise RuntimeError(f"PPTX extraction failed: {e}")
